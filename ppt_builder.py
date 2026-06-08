"""
10장 슬라이드(JSON) → .pptx (16:9 다크테마).
슬라이드 스키마: {title, points:[{head, desc}]}
본문은 '핵심포인트(굵게) + 부연설명(회색)' 블록으로 카드를 채운다.
"""
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

BG      = RGBColor(0x0F, 0x17, 0x2A)
TEXT    = RGBColor(0xF8, 0xFA, 0xFC)
MUTED   = RGBColor(0x94, 0xA3, 0xB8)
SUBTLE  = RGBColor(0xCB, 0xD5, 0xE1)
ACCENT  = RGBColor(0x60, 0xA5, 0xFA)
ACCENT2 = RGBColor(0xFB, 0xBF, 0x24)
SUCCESS = RGBColor(0x4A, 0xDE, 0x80)
SURFACE = RGBColor(0x1E, 0x29, 0x3B)
BORDER  = RGBColor(0x33, 0x41, 0x55)
FONT = "맑은 고딕"

# 한 줄에 들어가는 대략 글자 수(폭 기준, 폰트별)
CHARS_HEAD = 30
CHARS_SUB = 42


def _bg(slide, prs):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    s.fill.solid(); s.fill.fore_color.rgb = BG
    s.line.fill.background(); s.shadow.inherit = False
    return s


def _text(slide, left, top, width, height, lines, size=24, bold=False,
          color=TEXT, align=PP_ALIGN.LEFT, font=FONT, space=2):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, Emu(0))
    if isinstance(lines, str):
        lines = [lines]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space)
        run = p.add_run(); run.text = line
        run.font.size = Pt(size); run.font.bold = bold
        run.font.color.rgb = color; run.font.name = font
        rPr = run._r.get_or_add_rPr()
        ea = rPr.find(qn("a:ea"))
        if ea is None:
            ea = rPr.makeelement(qn("a:ea"), {"typeface": font}); rPr.append(ea)
        else:
            ea.set("typeface", font)
    return tb


def _box(slide, left, top, width, height, fill=SURFACE, line=BORDER):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = line; s.line.width = Pt(1)
    s.adjustments[0] = 0.05; s.shadow.inherit = False
    return s


def _chip(slide, left, top, text, fill=ACCENT, color=BG):
    w = Inches(0.15 * len(text) + 0.55)
    c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, Inches(0.42))
    c.fill.solid(); c.fill.fore_color.rgb = fill
    c.line.fill.background(); c.adjustments[0] = 0.5; c.shadow.inherit = False
    tf = c.text_frame
    for m in ("margin_top", "margin_bottom"):
        setattr(tf, m, Emu(20000))
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = color; r.font.name = FONT
    return c


def _bar(slide, left, top, width, color=ACCENT):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.08))
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background(); s.shadow.inherit = False
    return s


def _dot(slide, left, top, color=ACCENT2, d=0.15):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, Inches(d), Inches(d))
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background(); s.shadow.inherit = False
    return s


def _title_block(slide, idx, total, title, color=ACCENT, chip=None):
    _chip(slide, Inches(0.9), Inches(0.55), chip or f"{idx+1:02d} / {total:02d}", fill=color)
    _text(slide, Inches(0.9), Inches(1.2), Inches(11.5), Inches(1.0),
          title, size=30, bold=True, color=TEXT)
    _bar(slide, Inches(0.9), Inches(2.15), Inches(2.0), color=color)


def _points_card(slide, points, top=Inches(2.55), color=ACCENT2,
                 head_size=19, sub_size=14):
    """head(굵게) + subs(세부 불렛, 회색) 블록들. 내용량에 맞춰 카드 높이 자동."""
    left = Inches(0.9); width = Inches(11.5)
    top_in = top / Inches(1)            # 인치 실수값
    max_bottom = 7.5 - 0.4
    head_left = Inches(1.5); head_w = Inches(10.5)
    sub_left = Inches(1.95); sub_w = Inches(10.0)

    pts = [p for p in points[:3]
           if p.get("head", "").strip() or [s for s in p.get("subs", []) if s.strip()]]

    # 1) 필요한 내용 높이 계산
    content = top_in + 0.30
    plan = []
    for p in pts:
        head = p.get("head", "").strip()
        subs = [s.strip() for s in p.get("subs", []) if s.strip()][:3]
        h_lines = max(1, (len(head) // CHARS_HEAD) + 1)
        block = {"head": head, "subs": [], "h_lines": h_lines}
        content += 0.38 * h_lines + 0.06
        for sub in subs:
            s_lines = max(1, (len(sub) // CHARS_SUB) + 1)
            block["subs"].append((sub, s_lines))
            content += 0.30 * s_lines
        content += 0.20
        plan.append(block)
    content += 0.15  # 하단 패딩

    card_bottom = min(max_bottom, content)
    # 내용이 적으면 카드를 살짝 키워 너무 납작하지 않게(최소 2.4인치)
    card_h = max(card_bottom - top_in, 2.4)
    _box(slide, left, top, width, Inches(card_h))

    # 2) 렌더
    y = top + Inches(0.3)
    for block in plan:
        _dot(slide, Inches(1.05), y + Inches(0.06), color=color)
        _text(slide, head_left, y, head_w, Inches(0.45 * block["h_lines"]),
              block["head"], size=head_size, bold=True, color=TEXT, space=1)
        y += Inches(0.38 * block["h_lines"]) + Inches(0.06)
        for sub, s_lines in block["subs"]:
            _text(slide, sub_left, y, sub_w, Inches(0.34 * s_lines),
                  f"• {sub}", size=sub_size, bold=False, color=SUBTLE, space=1)
            y += Inches(0.30 * s_lines)
        y += Inches(0.20)


def build_pptx(slides: list[dict], source_url: str = "", audience: str = "") -> BytesIO:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    total = len(slides)

    for idx, sd in enumerate(slides):
        title = (sd.get("title") or "").strip() or f"슬라이드 {idx+1}"
        points = sd.get("points") or []
        slide = prs.slides.add_slide(blank)
        _bg(slide, prs)

        if idx == 0:
            # 표지
            chip_txt = (f"{audience} 강의자료" if audience else "AI 생성 강의자료 · 뉴스 브리핑")
            _chip(slide, Inches(0.9), Inches(1.4), chip_txt, fill=ACCENT2)
            _text(slide, Inches(0.9), Inches(2.4), Inches(11.5), Inches(2.4),
                  title, size=44, bold=True, color=TEXT)
            sub = points[0].get("head") if points else ""
            if sub:
                _text(slide, Inches(0.9), Inches(5.3), Inches(11.5), Inches(1.0),
                      sub, size=22, color=MUTED)
            if source_url:
                _text(slide, Inches(0.9), Inches(6.8), Inches(11.5), Inches(0.4),
                      f"출처: {source_url}", size=12, color=MUTED)

        elif idx == total - 1:
            # 정리
            _title_block(slide, idx, total, title, color=SUCCESS, chip="정리")
            _points_card(slide, points, color=SUCCESS)

        elif idx == total - 2:
            # 토론 질문 (head=질문, desc=보조)
            _title_block(slide, idx, total, title, color=ACCENT2, chip="수업 토론")
            _points_card(slide, points, color=ACCENT2, head_size=20)

        else:
            # 본문 — 7(전공)/8(강조)은 색 구분
            if "전공" in title:
                c = SUCCESS
            elif "강조" in title or idx == 7:
                c = ACCENT2
            else:
                c = ACCENT
            _title_block(slide, idx, total, title, color=c)
            _points_card(slide, points, color=c)

    out = BytesIO()
    prs.save(out)
    out.seek(0)
    return out
